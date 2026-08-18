from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------- Data ----------
reps = ['Abdullah', 'Mustafa', 'Aamir', 'Siraj', 'Kashif', 'Asif']
months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul']
monthly = {
    'Abdullah': [75000, 80000, 70000, 85000, 65000, 88000, 85000],
    'Mustafa':  [125000, 130000, 150000, 145000, 152000, 170000, 180000],
    'Aamir':    [25000, 35000, 60000, 45000, 50000, 52000, 58000],
    'Siraj':    [95000, 105000, 110000, 102000, 106000, 125000, 115000],
    'Kashif':   [35000, 33000, 32000, 36000, 33000, 32500, 38000],
    'Asif':     [46000, 39000, 47000, 48300, 50000, 49000, 53000],
}
ytd_target = {'Abdullah': 1200000, 'Mustafa': 1600000, 'Aamir': 1200000,
              'Siraj': 1000000, 'Kashif': 1350000, 'Asif': 1200000}
ytd_sales = {'Abdullah': 1748000, 'Mustafa': 2652000, 'Aamir': 1525000,
             'Siraj': 1758000, 'Kashif': 1589500, 'Asif': 1532300}
july = {r: monthly[r][6] for r in reps}
forecast = {r: int(july[r] * 1.2) for r in reps}
july_total = sum(july.values())
forecast_total = int(july_total * 1.2)

chart_dir = r'C:\Users\DELL LATITUDE 5520\.openclaw\workspace\charts'
out_path = r'C:\Users\DELL LATITUDE 5520\.openclaw\workspace\Sales_Performance_and_Plan.pptx'

# ---------- Theme ----------
NAVY = RGBColor(0x1F, 0x3F, 0x6E)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT = RGBColor(0xDE, 0xEB, 0xF7)
DARK = RGBColor(0x26, 0x26, 0x26)
GREY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x3B, 0x8E, 0x41)
ORANGE = RGBColor(0xE8, 0x8A, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]

def add_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, line=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line:
        sh.line.color.rgb = color
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def add_text(slide, l, t, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return tb

def add_bullets(slide, l, t, w, h, items, size=16, color=DARK, gap=6):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            txt, lvl, bold = it
        else:
            txt, lvl, bold = it, 0, False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = txt
        p.level = lvl
        p.space_after = Pt(gap)
        r = p.runs[0]
        r.font.size = Pt(size if lvl == 0 else size-2)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = 'Calibri'
    return tb

def add_title_bar(slide, title, subtitle=None, num=None):
    add_rect(slide, 0, 0, SW, Inches(1.0), NAVY)
    add_rect(slide, 0, Inches(1.0), SW, Inches(0.06), ORANGE)
    add_text(slide, Inches(0.5), Inches(0.12), SW - Inches(1.0), Inches(0.5),
             title, size=26, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.58), SW - Inches(1.0), Inches(0.4),
                 subtitle, size=13, color=LIGHT)
    if num:
        add_text(slide, SW - Inches(1.2), Inches(0.12), Inches(0.9), Inches(0.7),
                 num, size=30, color=RGBColor(0x9D, 0xC3, 0xE6), bold=True, align=PP_ALIGN.RIGHT)

def add_picture_center(slide, path, l, t, w, h):
    pic = slide.shapes.add_picture(path, l, t, w, h)
    return pic

# ================= SLIDE 1: TITLE =================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_rect(s, 0, Inches(2.2), SW, Inches(0.08), ORANGE)
add_text(s, Inches(0.9), Inches(1.0), SW - Inches(1.8), Inches(1.2),
         'Sales Performance Review', size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(2.5), SW - Inches(1.8), Inches(0.8),
         '& Next Month Sales Plan', size=30, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(4.3), SW - Inches(1.8), Inches(0.5),
         'Sales Team Meeting  |  support@stellarstech.com', size=16, color=RGBColor(0x9D, 0xC3, 0xE6), align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(4.9), SW - Inches(1.8), Inches(0.5),
         'August 2026', size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# ================= SLIDE 2: AGENDA =================
s = prs.slides.add_slide(blank)
add_bg(s)
add_title_bar(s, 'Agenda', 'Meeting flow', '02')
agenda = [
    '1.  Performance Overview — key numbers',
    '2.  YTD Sales vs Targets',
    '3.  Monthly Sales Trend',
    '4.  Rep Performance & Sales Share',
    '5.  Next Month Sales Plan — key components',
    '6.  Sales Forecast — +20% growth',
    '7.  How to Boost Sales — team guide',
    '8.  Action Items & Next Steps',
]
add_bullets(s, Inches(1.0), Inches(1.6), Inches(11.3), Inches(5.0), agenda, size=20, gap=14)

# ================= SLIDE 3: KEY HIGHLIGHTS =================
s = prs.slides.add_slide(blank)
add_bg(s)
add_title_bar(s, 'Performance Overview — Key Highlights', 'Summary of YTD results', '03')

ytd_sales_total = sum(ytd_sales.values())
cards = [
    ('Total YTD Sales', f'{ytd_sales_total / 1e6:.2f}M', GREEN),
    ('Reps Above Target', '6 / 6', GREEN),
    ('Growth Jan → Jul', '+32%', BLUE),
    ('Top Performer', 'Siraj (175.8%)', ORANGE),
]
cw, ch, gapx = Inches(2.7), Inches(2.0), Inches(0.35)
start_x = Inches(0.55)
for i, (label, val, color) in enumerate(cards):
    x = start_x + i * (cw + gapx)
    y = Inches(1.7)
    add_rect(s, x, y, cw, ch, LIGHT)
    add_rect(s, x, y, cw, Inches(0.12), color)
    add_text(s, x + Inches(0.15), y + Inches(0.35), cw - Inches(0.3), Inches(0.5), label, size=13, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y + Inches(0.85), cw - Inches(0.3), Inches(0.8), val, size=30, color=color, bold=True, align=PP_ALIGN.CENTER)

hl = [
    'All 6 sales reps exceeded their YTD targets — no one is below 100%.',
    'Mustafa leads in absolute volume (2,652,000); Siraj leads in efficiency (175.8% of target).',
    'Monthly trend is consistently rising: 401K in Jan to 529K in July.',
    'Next step: build on this momentum with a +20% growth plan for next month.',
]
add_bullets(s, Inches(0.55), Inches(4.15), Inches(12.2), Inches(2.8), hl, size=16, gap=8)

# ================= SLIDE 4: YTD SALES VS TARGET =================
s = prs.slides.add_slide(blank)
add_bg(s)
add_title_bar(s, 'YTD Sales vs Target by Rep', 'Every rep exceeded target', '04')
add_picture_center(s, os.path.join(chart_dir, '2_ytd_sales_vs_target.png'),
                   Inches(1.7), Inches(1.4), Inches(9.9), Inches(5.6))

# ================= SLIDE 5: MONTHLY TREND =================
s = prs.slides.add_slide(blank)
add_bg(s)
add_title_bar(s, 'Monthly Sales Trend (Jan – Jul)', 'Consistent upward momentum', '05')
add_picture_center(s, os.path.join(chart_dir, '3_monthly_trend.png'),
                   Inches(2.1), Inches(1.5), Inches(9.1), Inches(5.1))
add_text(s, Inches(2.1), Inches(6.6), Inches(9.1), Inches(0.5),
         'Total monthly sales grew from 401K to 529K (+32%)', size=14, color=GREY, align=PP_ALIGN.CENTER)

# ================= SLIDE 6: REP PERFORMANCE =================
s = prs.slides.add_slide(blank)
add_bg(s)
add_title_bar(s, 'Rep Performance — % of Target', 'Siraj & Mustafa lead the team', '06')
add_picture_center(s, os.path.join(chart_dir, '5_pct_target_achieved.png'),
                   Inches(2.1), Inches(1.4), Inches(9.1), Inches(5.2))

# ================= SLIDE 7: SALES SHARE =================
s = prs.slides.add_slide(blank)
add_bg(s)
add_title_bar(s, 'YTD Sales Share by Rep', 'Contribution split', '07')
add_picture_center(s, os.path.join(chart_dir, '4_ytd_sales_share_pie.png'),
                   Inches(3.6), Inches(1.3), Inches(6.1), Inches(5.6))

# ================= SLIDE 8: NEXT MONTH SALES PLAN =================
s = prs.slides.add_slide(blank)
add_bg(s)
add_title_bar(s, 'Next Month Sales Plan — Key Components', 'A structured, measurable roadmap', '08')

plan_left = [
    ('1.  Lead Generation', 0, True),
    ('   •  Increase daily prospecting calls/emails by 25%', 1, False),
    ('   •  Reactivate old/pending leads database', 1, False),
    ('2.  Customer Retention', 0, True),
    ('   •  Follow-up every existing account at least twice', 1, False),
    ('   •  Resolve pending queries within 24 hours', 1, False),
    ('3.  Upselling & Cross-selling', 0, True),
    ('   •  Offer add-on products to top 20 customers', 1, False),
    ('   •  Bundle offers to increase average order value', 1, False),
]
plan_right = [
    ('4.  Clear Targets', 0, True),
    ('   •  Individual target set per rep (see forecast)', 1, False),
    ('   •  Weekly milestones to track progress', 1, False),
    ('5.  Daily Tracking & Reporting', 0, True),
    ('   •  Daily activity log + CRM updates', 1, False),
    ('   •  Weekly review meeting every Monday', 1, False),
    ('6.  Incentives', 0, True),
    ('   •  Top performer bonus & recognition', 1, False),
    ('   •  Team reward for hitting +20% target', 1, False),
]
add_bullets(s, Inches(0.55), Inches(1.5), Inches(6.1), Inches(5.6), plan_left, size=15, gap=7)
add_bullets(s, Inches(6.9), Inches(1.5), Inches(6.0), Inches(5.6), plan_right, size=15, gap=7)

# ================= SLIDE 9: FORECAST +20% =================
s = prs.slides.add_slide(blank)
add_bg(s)
add_title_bar(s, 'Sales Forecast — Next Month (+20%)', 'Baseline: July actuals', '09')

rows = len(reps) + 2
cols = 4
tbl_shape = s.shapes.add_table(rows, cols, Inches(1.1), Inches(1.5), Inches(11.1), Inches(4.6))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(3.3)
tbl.columns[1].width = Inches(2.6)
tbl.columns[2].width = Inches(2.6)
tbl.columns[3].width = Inches(2.6)

headers = ['Sales Rep', 'July Sales', 'Forecast (+20%)', 'Growth Amount']
data_rows = []
for r in reps:
    data_rows.append([r, f'{july[r]:,}', f'{forecast[r]:,}', f'{forecast[r]-july[r]:,}'])
data_rows.append(['TOTAL', f'{july_total:,}', f'{forecast_total:,}', f'{forecast_total-july_total:,}'])

for c, h in enumerate(headers):
    cell = tbl.cell(0, c)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.runs[0].font.bold = True; p.runs[0].font.size = Pt(15); p.runs[0].font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

for ri, row in enumerate(data_rows, start=1):
    is_total = (row[0] == 'TOTAL')
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = val
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(14)
        p.runs[0].font.bold = is_total
        if ci == 0:
            p.alignment = PP_ALIGN.LEFT
        else:
            p.alignment = PP_ALIGN.CENTER
        if is_total:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
        elif ri % 2 == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF6, 0xFA)
        if ci == 2:
            p.runs[0].font.color.rgb = GREEN

add_text(s, Inches(1.1), Inches(6.4), Inches(11.1), Inches(0.5),
         'Next month total target: ' + f'{forecast_total:,}' + '  (from ' + f'{july_total:,}' + ' baseline)',
         size=15, color=DARK, bold=True, align=PP_ALIGN.CENTER)

# ================= SLIDE 10: HOW TO BOOST SALES =================
s = prs.slides.add_slide(blank)
add_bg(s)
add_title_bar(s, 'How to Boost Sales — Team Guide', 'Actionable tactics for every rep', '10')

boost = [
    ('1.  Prospect Daily', 0, True),
    ('   •  Set a minimum of 15 new prospect touch-points every day', 1, False),
    ('2.  Improve Conversion', 0, True),
    ('   •  Speed up quote turnaround to under 4 hours', 1, False),
    ('   •  Follow up hot leads within 30 minutes', 1, False),
    ('3.  Grow Account Value', 0, True),
    ('   •  Upsell/cross-sell to existing customers first', 1, False),
    ('   •  Ask every customer for one referral', 1, False),
    ('4.  Track & Learn', 0, True),
    ('   •  Log every call, quote and outcome in CRM daily', 1, False),
    ('   •  Review your own win/loss each week', 1, False),
]
add_bullets(s, Inches(0.55), Inches(1.4), Inches(6.0), Inches(5.6), boost, size=16, gap=9)

boost_right = [
    ('5.  Focus on High-Value Accounts', 0, True),
    ('   •  Prioritize the top 20% customers by value', 1, False),
    ('6.  Team Collaboration', 0, True),
    ('   •  Share winning pitches & objections handled', 1, False),
    ('   •  Pair junior reps with top performers', 1, False),
    ('7.  Stay Consistent', 0, True),
    ('   •  Hit daily activity targets every day, not just at month-end', 1, False),
    ('8.  Use the Data', 0, True),
    ('   •  Revisit this dashboard weekly to spot trends', 1, False),
]
add_bullets(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(5.6), boost_right, size=16, gap=9)

# ================= SLIDE 11: ACTION ITEMS =================
s = prs.slides.add_slide(blank)
add_bg(s)
add_title_bar(s, 'Action Items & Next Steps', 'Owners and timing', '11')

actions = [
    ('Action', 'Owner', 'When'),
    ('Assign individual +20% targets', 'Sales Manager', 'Today'),
    ('Launch daily activity tracker', 'Team Lead', 'This week'),
    ('Reactivate old leads list', 'All reps', 'This week'),
    ('Weekly review meeting (Mon)', 'Sales Manager', 'Weekly'),
    ('Top-performer incentive plan', 'Management', 'This month'),
]
rows = len(actions)
tbl2 = s.shapes.add_table(rows, 3, Inches(1.1), Inches(1.6), Inches(11.1), Inches(4.0)).table
tbl2.columns[0].width = Inches(6.2)
tbl2.columns[1].width = Inches(2.6)
tbl2.columns[2].width = Inches(2.3)
for ri, row in enumerate(actions):
    for ci, val in enumerate(row):
        cell = tbl2.cell(ri, ci)
        cell.text = val
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(14)
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            p.runs[0].font.bold = True; p.runs[0].font.color.rgb = WHITE
        elif ri % 2 == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF6, 0xFA)
        if ci == 0:
            p.alignment = PP_ALIGN.LEFT
        else:
            p.alignment = PP_ALIGN.CENTER

# ================= SLIDE 12: THANK YOU =================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_rect(s, 0, Inches(3.0), SW, Inches(0.08), ORANGE)
add_text(s, Inches(0.9), Inches(2.2), SW - Inches(1.8), Inches(1.0),
         'Thank You', size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(3.4), SW - Inches(1.8), Inches(0.6),
         'Let\'s hit the +20% target together', size=22, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(4.4), SW - Inches(1.8), Inches(0.5),
         'support@stellarstech.com', size=14, color=RGBColor(0x9D, 0xC3, 0xE6), align=PP_ALIGN.CENTER)

prs.save(out_path)
print('SAVED', out_path)
print('july_total', july_total, 'forecast_total', forecast_total)
